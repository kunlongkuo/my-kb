import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd

# Supported file extensions by category
TEXT_EXTENSIONS = ['.txt', '.md', '.py', '.js', '.css', '.html', '.json', '.yaml', '.yml', '.csv', '.log']
IMAGE_EXTENSIONS = ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp']
MEDIA_EXTENSIONS = ['.mp4', '.mov', '.mp3', '.wav', '.avi', '.mkv']

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()

def get_file_type(file_path):
    """Return the file type category: 'text', 'image', 'media', or 'unsupported'."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in TEXT_EXTENSIONS:
        return 'text'
    elif ext in ['.pdf']:
        return 'pdf'
    elif ext in ['.docx']:
        return 'docx'
    elif ext in ['.pptx']:
        return 'pptx'
    elif ext in ['.xlsx', '.xls']:
        return 'excel'
    elif ext in IMAGE_EXTENSIONS:
        return 'image'
    elif ext in MEDIA_EXTENSIONS:
        return 'media'
    else:
        return 'unsupported'

def process_file(file_path):
    """
    Process a file and extract its text content.
    Returns (text_content, file_type) tuple.
    For images, text_content will be None (handled separately by vector_store).
    """
    file_type = get_file_type(file_path)
    
    try:
        if file_type == 'text':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(), 'text'
        elif file_type == 'pdf':
            return extract_text_from_pdf(file_path), 'pdf'
        elif file_type == 'docx':
            return extract_text_from_docx(file_path), 'docx'
        elif file_type == 'pptx':
            return extract_text_from_pptx(file_path), 'pptx'
        elif file_type == 'excel':
            return extract_text_from_xlsx(file_path), 'excel'
        elif file_type == 'image':
            return None, 'image'  # Images handled by vector_store.add_image_document
        elif file_type == 'media':
            return None, 'media'  # Skip media for now
        else:
            return None, 'unsupported'
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return None, 'error'
